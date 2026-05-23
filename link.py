# -*- coding: utf-8 -*-
import sys, os, shutil
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
desktop = os.path.join(os.path.expanduser("~"), "Desktop")
src = os.path.join(desktop, "Second-Hand_10页精炼版.pptx")
tmp = os.path.join(desktop, "_temp.pptx")
shutil.copy2(src, tmp)
p = Presentation(tmp)
s = p.slides[8]
for sh in s.shapes:
    if not sh.has_text_frame: continue
    for pf in sh.text_frame.paragraphs:
        if 'freevian.top' in pf.text:
            for r in pf.runs: r.text = ''
            run = pf.add_run()
            run.text = 'http://freevian.top:5000'
            run.font.size = Pt(24)
            run.hyperlink.address = 'http://freevian.top:5000'
            print('OK')
            break
    else: continue
    break
p.save(tmp)
shutil.copy2(tmp, src)
os.remove(tmp)
print('DONE')
